"""Generate COOK demo deck -> docs/cook-demo.pptx

Direction: "Proof, Not Vibes" — a measurement story. Same Pass design
system (paper, ink, flame, Georgia italic display, approval stamp) with
one new device: oversized proof numerals wherever the deck makes a claim.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

PAPER = RGBColor(0xF7, 0xF2, 0xE9)
INK = RGBColor(0x1B, 0x17, 0x12)
FLAME = RGBColor(0xD4, 0x4A, 0x24)
MUTED = RGBColor(0x6E, 0x64, 0x55)
CARD = RGBColor(0xFF, 0xFC, 0xF6)
LINE = RGBColor(0xD9, 0xCF, 0xBE)

DISPLAY = "Georgia"
BODY = "Segoe UI"
ELL = "\u2026"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
SW, SH = prs.slide_width, prs.slide_height


def slide():
    s = prs.slides.add_slide(BLANK)
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    bg.fill.solid()
    bg.fill.fore_color.rgb = PAPER
    bg.line.fill.background()
    bg.shadow.inherit = False
    return s


def _set_run(p, text, size, color, bold=False, italic=False, font=BODY, spacing=None):
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.color.rgb = color
    r.font.bold = bold
    r.font.italic = italic
    r.font.name = font
    if spacing is not None:
        r.font._rPr.set("spc", str(spacing))
    return r


def text(s, l, t, w, h, lines, anchor=MSO_ANCHOR.TOP):
    box = s.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    for i, spec in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = spec.get("align", PP_ALIGN.LEFT)
        p.space_after = Pt(spec.get("after", 6))
        p.space_before = Pt(spec.get("before", 0))
        if spec.get("line"):
            p.line_spacing = spec["line"]
        for run in spec["runs"]:
            _set_run(p, *run[:5], **(run[5] if len(run) > 5 else {}))
    return box


def R(t, size, color, bold=False, italic=False, font=BODY, spacing=None):
    extra = {"font": font}
    if spacing is not None:
        extra["spacing"] = spacing
    return (t, size, color, bold, italic, extra)


def hairline(s, l, t, w):
    ln = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, Pt(1))
    ln.fill.solid()
    ln.fill.fore_color.rgb = LINE
    ln.line.fill.background()
    ln.shadow.inherit = False
    return ln


def furniture(s, section, n, total=8):
    hairline(s, Inches(0.9), Inches(0.62), Inches(11.53))
    text(s, Inches(0.9), Inches(0.74), Inches(9), Inches(0.35), [
        {"runs": [
            R("COOK", 11, INK, bold=True, spacing=300),
            R("   ·   " + section.upper(), 11, MUTED, bold=True, spacing=300),
        ]}
    ])
    text(s, SW - Inches(2.0), Inches(0.74), Inches(1.1), Inches(0.35), [
        {"runs": [R(f"{n:02d} / {total:02d}", 11, MUTED, bold=True, spacing=200)], "align": PP_ALIGN.RIGHT}
    ])


def headline(s, runs, t=Inches(1.3), size=52, w=Inches(11.53)):
    if isinstance(runs, str):
        runs = [R(runs, size, INK, italic=True, font=DISPLAY)]
    text(s, Inches(0.9), t, w, Inches(1.7), [{"runs": runs, "line": 1.02}])


def stamp(s, l, t, label="APPROVAL REQUIRED", rot=-8, w=Inches(2.5), h=Inches(0.62), size=13):
    chip = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    chip.adjustments[0] = 0.5
    chip.fill.background()
    chip.line.color.rgb = FLAME
    chip.line.width = Pt(1.75)
    chip.shadow.inherit = False
    chip.rotation = rot
    tf = chip.text_frame
    tf.word_wrap = False
    tf.margin_left = tf.margin_right = Inches(0.05)
    tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    _set_run(p, label, size, FLAME, bold=True, font=BODY, spacing=300)
    return chip


def card(s, l, t, w, h):
    c = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    c.adjustments[0] = 0.045
    c.fill.solid()
    c.fill.fore_color.rgb = CARD
    c.line.color.rgb = LINE
    c.line.width = Pt(1)
    c.shadow.inherit = False
    return c


def bignum(s, l, t, num, label, num_size=64, w=Inches(3.4), num_color=FLAME):
    text(s, l, t, w, Inches(1.1), [
        {"runs": [R(num, num_size, num_color, bold=True, font=DISPLAY, spacing=-100)], "line": 0.95}
    ])
    text(s, l, t + Inches(num_size / 72.0 + 0.12), w, Inches(0.9), [
        {"runs": [R(label, 13.5, MUTED)], "line": 1.15}
    ])


# ============ 1. TITLE ============
s = slide()
hairline(s, Inches(0.9), Inches(0.62), Inches(11.53))
text(s, Inches(0.9), Inches(0.74), Inches(9), Inches(0.35), [
    {"runs": [R("CLUB EVENT OPS", 11, MUTED, bold=True, spacing=500)]}
])
text(s, Inches(0.82), Inches(1.5), Inches(11.8), Inches(2.7), [
    {"runs": [R("COOK", 180, INK, bold=True, font=DISPLAY, spacing=-200)], "line": 0.9}
])
text(s, Inches(0.9), Inches(4.15), Inches(11.5), Inches(1.0), [
    {"runs": [
        R("Club events, ", 30, INK, italic=True, font=DISPLAY),
        R("measured.", 30, FLAME, italic=True, font=DISPLAY),
    ]}
])
text(s, Inches(0.9), Inches(5.05), Inches(10.8), Inches(1.2), [
    {"runs": [R("COOK runs the event & keeps the score — turnout, spend, replies, follow-through —", 17, MUTED)], "after": 2},
    {"runs": [R("so every event makes the next one sharper.", 17, MUTED)]},
])
hairline(s, Inches(0.9), Inches(6.7), Inches(11.53))
text(s, Inches(0.9), Inches(6.85), Inches(11.53), Inches(0.4), [
    {"runs": [
        R("DEMO OVERVIEW", 11, MUTED, bold=True, spacing=400),
        R("      8 SLIDES", 11, MUTED, bold=True, spacing=400),
    ]}
])

# ============ 2. THE PROBLEM ============
s = slide()
furniture(s, "The Problem", 2)
headline(s, "Nobody Keeps the Score")
rows = [
    ("Nobody knows what worked.", "The event ends, the group chat moves on, & the spreadsheet of who showed up never happens."),
    ("The next event starts from zero.", "No attendance history, no cost records, no follow-up list — every officer re-learns the same lessons."),
    ("AI tools don’t help.", "They chat about the work, or send & spend without asking. None of them count anything."),
]
top = Inches(2.5)
for i, (h, b) in enumerate(rows):
    y = top + Inches(1.42) * i
    text(s, Inches(0.9), y, Inches(0.6), Inches(0.5), [
        {"runs": [R("—", 22, FLAME, bold=True)]}
    ])
    text(s, Inches(1.55), y, Inches(10.8), Inches(1.3), [
        {"runs": [R(h, 21, INK, bold=True)], "after": 2},
        {"runs": [R(b, 15, MUTED)], "line": 1.15},
    ])

# ============ 3. WHAT COOK IS ============
s = slide()
furniture(s, "What COOK Is", 3)
headline(s, "Run the Event. Keep the Score.", size=46)
cards = [
    ("The Pass", "Every event, task, run-of-show & partner in one place — scoped to your club, shared by your officers."),
    ("The Kitchen Staff", "AI skills prep the work: task lists, outreach drafts, supply carts, food runs. You stay on the pass."),
    ("The Scorecard", "Every event records its numbers — turnout, spend, replies, follow-through — and feeds them forward."),
]
cw, ch = Inches(3.68), Inches(3.1)
gap = Inches(0.25)
top = Inches(2.55)
for i, (h, b) in enumerate(cards):
    l = Inches(0.9) + (cw + gap) * i
    card(s, l, top, cw, ch)
    text(s, l + Inches(0.32), top + Inches(0.35), cw - Inches(0.64), ch - Inches(0.7), [
        {"runs": [R(h, 22, FLAME, italic=True, font=DISPLAY)], "after": 10},
        {"runs": [R(b, 14.5, INK)], "line": 1.25},
    ])
text(s, Inches(0.9), Inches(6.2), Inches(11.53), Inches(0.8), [
    {"runs": [
        R("Not a social app. Not five chatbots. ", 16, MUTED),
        R("One spine, skills on the side, numbers on the record.", 16, INK, bold=True),
    ]}
])

# ============ 4. WHAT GETS MEASURED ============
s = slide()
furniture(s, "What Gets Measured", 4)
headline(s, "What Gets Measured")
stats = [
    ("Show-up rate", "RSVPs vs. actual check-ins — the number every club argues about, settled."),
    ("Cost per head", "Supplies & food, divided by who came. Budgets stop being guesses."),
    ("Reply rate", "Outreach answered vs. sent — which partners & channels actually respond."),
    ("Follow-through", "Post-event tasks closed vs. opened. The debrief that actually happens."),
]
cw, ch = Inches(5.63), Inches(1.62)
gap = Inches(0.27)
top = Inches(2.6)
for i, (h, b) in enumerate(cards := stats):
    l = Inches(0.9) + (cw + gap) * (i % 2)
    t = top + (ch + gap) * (i // 2)
    card(s, l, t, cw, ch)
    text(s, l + Inches(0.32), t + Inches(0.24), cw - Inches(0.64), ch - Inches(0.44), [
        {"runs": [R(h, 18, INK, bold=True)], "after": 4},
        {"runs": [R(b, 13.5, MUTED)], "line": 1.2},
    ])
text(s, Inches(0.9), Inches(6.35), Inches(11.53), Inches(0.6), [
    {"runs": [
        R("No PII in the numbers. ", 15, INK, bold=True),
        R("Counts & rates, never names — attendance, not surveillance.", 15, MUTED),
    ]}
])

# ============ 5. THE SKILLS ============
s = slide()
furniture(s, "The Skills", 5)
headline(s, "Every Station Feeds the Score", size=46)
rows = [
    ("Plan Club Event", "Sets targets up front: headcount, budget, capacity. You can’t beat a number you never set."),
    ("Luma Event Invite", "The draft invite is where RSVP tracking starts — signups flow back into the record."),
    ("Amazon Event Supplies", "500+ reviews, 4.6+ stars, arrives 2+ days early — quality bars you can audit, with the spend logged."),
    ("Event Food Order", "Order sheet, runner logistics, receipt — cost per head falls out of the paperwork."),
    ("Metrics & Follow-Ups", "Reads the event’s records after the fact: what hit target, what slipped, who to thank."),
    ("Personal Daily Brief", "Each morning: what’s on, what’s waiting on you, what’s overdue — your own follow-through, counted."),
]
top = Inches(2.35)
col_w = Inches(5.63)
for i, (h, b) in enumerate(rows):
    l = Inches(0.9) + (col_w + Inches(0.27)) * (i % 2)
    t = top + Inches(1.52) * (i // 2)
    text(s, l, t, col_w, Inches(1.45), [
        {"runs": [R("—  ", 16, FLAME, bold=True), R(h, 16.5, INK, bold=True)], "after": 3},
        {"runs": [R(b, 12.5, MUTED)], "line": 1.18},
    ])

# ============ 6. THE LOOP ============
s = slide()
furniture(s, "The Loop", 6)
headline(s, "Every Event Sharpens the Next", size=46)
steps = [
    ("Set targets", "Headcount, budget & capacity lock in at the plan stage — before anyone spends a dollar."),
    ("Run the night", "Check-ins, spend & outreach land in the record as the event happens, not after from memory."),
    ("Read the score", "Show-up rate, cost per head, reply rate — against the targets you set in step 1."),
    ("Feed it forward", "Follow-ups go out, lessons attach to the next plan, targets adjust. The loop closes."),
]
top = Inches(2.55)
row_h = Inches(1.08)
for i, (h, b) in enumerate(steps):
    t = top + row_h * i
    text(s, Inches(0.9), t - Inches(0.12), Inches(1.1), Inches(0.9), [
        {"runs": [R(f"{i+1:02d}", 34, FLAME, italic=True, font=DISPLAY)]}
    ])
    text(s, Inches(2.1), t, Inches(9.4), Inches(1.0), [
        {"runs": [R(h, 19, INK, bold=True)], "after": 2},
        {"runs": [R(b, 13.5, MUTED)], "line": 1.15},
    ])
    if i < 3:
        hairline(s, Inches(2.1), t + row_h - Inches(0.14), Inches(10.33))
stamp(s, Inches(10.55), Inches(1.15), rot=7, w=Inches(2.15), h=Inches(0.56), size=11)

# ============ 7. SAFETY ============
s = slide()
furniture(s, "Safety", 7)
headline(s, [R("Nothing Leaves the Kitchen", 46, INK, italic=True, font=DISPLAY)], size=46)
text(s, Inches(0.9), Inches(1.95), Inches(11.53), Inches(0.9), [
    {"runs": [R("without your stamp.", 46, FLAME, italic=True, font=DISPLAY)]}
])
rows = [
    ("Approve-before-send, everywhere", "invites, outreach, orders & edits each pause for an explicit yes — per action, not per session."),
    ("Dry-run safe by default", "demo mode builds full carts & drafts without ever paying, publishing, or sending."),
    ("Every number shows its source", "each metric traces to the record it came from — verify any claim in two clicks."),
    ("Counts, not people", "attendance & rates only. No names, no message bodies, no PII in the numbers."),
]
top = Inches(2.95)
for i, (h, b) in enumerate(rows):
    t = top + Inches(0.92) * i
    text(s, Inches(0.9), t, Inches(0.5), Inches(0.5), [
        {"runs": [R("—", 20, FLAME, bold=True)]}
    ])
    text(s, Inches(1.55), t, Inches(8.6), Inches(0.9), [
        {"runs": [R(h + " — ", 15.5, INK, bold=True), R(b, 13.5, MUTED)], "line": 1.15},
    ])
stamp(s, Inches(10.35), Inches(3.65), rot=-9, w=Inches(2.5), h=Inches(0.62), size=12)
stamp(s, Inches(10.7), Inches(4.8), rot=5, w=Inches(2.1), h=Inches(0.56), size=11, label="HUMAN SAYS YES")

# ============ 8. WHERE WE ARE ============
s = slide()
furniture(s, "Where We Are", 8)
headline(s, "Keeping Score Tonight")
bignum(s, Inches(0.9), Inches(2.45), "8", "skills written & ready — plan, invite, supplies, food, metrics, brief & more", num_size=72)
bignum(s, Inches(4.85), Inches(2.45), "1", "queue every draft, cart & invite must pass through before anything happens", num_size=72)
bignum(s, Inches(8.8), Inches(2.45), "0", "messages sent, orders placed, or invites published without a human yes", num_size=72)
hairline(s, Inches(0.9), Inches(4.6), Inches(11.53))
cols = [
    ("On the Pass", INK, "Full event management, the Approvals queue & a seeded demo night — ready to walk through now."),
    ("Still Prepping", MUTED, "AI skills answer with placeholders today; sign-in is a demo shortcut; approved drafts don’t send yet."),
    ("Next Up", FLAME, "A real model behind the skills, real accounts & per-club permissions, and live sending on approval."),
]
cw = Inches(3.68)
gap = Inches(0.25)
top = Inches(4.85)
for i, (h, hc, b) in enumerate(cols):
    l = Inches(0.9) + (cw + gap) * i
    text(s, l, top, cw, Inches(1.7), [
        {"runs": [R(h, 18, hc, italic=True, font=DISPLAY)], "after": 6},
        {"runs": [R(b, 12.5, INK if i != 1 else MUTED)], "line": 1.2},
    ])
text(s, Inches(0.9), Inches(6.6), Inches(11.53), Inches(0.6), [
    {"runs": [
        R("The scoreboard is built. ", 17, INK, italic=True, font=DISPLAY),
        R("We’re just firing the burners" + ELL, 17, FLAME, italic=True, font=DISPLAY),
    ], "align": PP_ALIGN.CENTER}
])

prs.save(r"C:\Users\bootcamp\projects\cursorHackcon\docs\cook-demo.pptx")
print("Saved docs/cook-demo.pptx —", len(prs.slides._sldIdLst), "slides")
